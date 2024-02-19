# from pptx import Presentation
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from openpyxl.drawing.image import Image as ExcelImage
import openpyxl


##############엑셓 파일 불러오는거######################

def read_excel_data(file_path):
    # 엑셀 파일을 읽어옵니다.
    df = pd.read_excel(file_path)

    # 필요한 열만 선택하여 데이터를 가공합니다.
    selected_columns = ['이름','이메일','사진',	'스펙',	'문항 1','문항 2','문항 3',	'answer 0',	'answer 1',	'answer 2',	'total','status']
    data = df[selected_columns].to_dict(orient='records')

    return data


# 엑셀 파일에서 데이터를 읽어와서 ppt를 생성하는 함수
def create_ppt_from_excel(path):
    result_path = path + '/result/result.xlsx'
    report_path = path + '/result/report.pptx'
    # 엑셀 파일에서 데이터를 읽어옵니다.
    data = read_excel_data(result_path)

    # 프레젠테이션 객체 생성
    prs = Presentation()
    
    # 각 행을 반복하면서 슬라이드 생성
    for row in data:
        name = row['이름']
        email = row['이메일']
        image_path = row['사진']
        spec = row['스펙']
        essay1 = row['answer 0']
        essay2 = row['answer 1']
        essay3 = row['answer 2']
        pass_fail = row['status']

        # 슬라이드 생성
        slide = prs.slides.add_slide(prs.slide_layouts[3])

        # 제목 설정
        slide.shapes.title.text = name
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.bold = True


        #ppt 구성 만들기
        body_shape = slide.placeholders[1]
        body_shape.width = Cm(24.4)
        body_shape.height = Cm(8.07)
        body_shape.left = Cm(0.5)
        body_shape.top = Cm(10.24)
        body_text = body_shape.text_frame
        body_p = body_text.add_paragraph()
        body_p.text = '평가 결과'
        body_p.font.size = Pt(18)
        body_p.font.bold = True

        
        
        # 이메일 추가 (제목 상자 오른쪽 아래에 추가)
        if email:
            title_shape = slide.shapes.title
            email_left = title_shape.left + title_shape.width - prs.slide_width * 0.28 #여기 조절로 이멜 오른쪽 아래 배치
            email_top = title_shape.top + title_shape.height
            email_shape = slide.shapes.add_textbox(email_left, email_top,
                                                prs.slide_width * 0.4, prs.slide_height * 0.1)
            email_shape.text = email
        

        # 이미지 추가
        if image_path:
            left = Cm(2.5)
            top = Cm(4.1) #이미지 아래쪽 이동
            width = Cm(7.62)
            height = Cm(6.73) #이미지 크기 조절
            slide.shapes.add_picture(image_path, left, top, width, height)

        if spec:
            spec_shape = slide.placeholders[2]
            spec_shape.width = Cm(12)
            spec_shape.height = Cm(3.62)
            spec_shape.left = Cm(12.7)
            spec_shape.top = Cm(6.69)
            spec_text = spec_shape.text_frame
            spec_p = spec_text.add_paragraph()
            spec_p.text = spec
            spec_p.font.size = Pt(18)


        if pass_fail:
            pass_fail_width = pass_fail_height = Cm(5)
            pass_fail_left = Cm(15.96)
            pass_fail_top = Cm(4.75)
            pass_fail_shape = slide.shapes.add_textbox(pass_fail_left, pass_fail_top, pass_fail_width, pass_fail_height)
            pass_fail_text = pass_fail_shape.text_frame
            pass_fail_p = pass_fail_text.add_paragraph()
            pass_fail_p.text = pass_fail
            pass_fail_p.font.size = Pt(32)
            pass_fail_p.font.bold = True

            
        if essay1:
            essay1_p = body_text.add_paragraph()
            essay1_p.text = '문항 1 점수'
            essay1_p.level = 1
            essay1_p.font.size = Pt(15)
            essay1_p = body_text.add_paragraph()
            essay1_p.text = essay1
            essay1_p.level= 1
            essay1_p.font.size = Pt(10)


        if essay2:
            essay2_p = body_text.add_paragraph()
            essay2_p.text = '문항 2 점수'
            essay2_p.level = 1
            essay2_p.font.size = Pt(15)
            essay2_p = body_text.add_paragraph()
            essay2_p.text = essay2
            essay2_p.level= 1
            essay2_p.font.size = Pt(10)

        if essay3:
            essay3_p = body_text.add_paragraph()
            essay3_p.text = '문항 3 점수'
            essay3_p.level = 1
            essay3_p.font.size = Pt(15)
            essay3_p = body_text.add_paragraph()
            essay3_p.text = essay3
            essay3_p.level= 1
            essay3_p.font.size = Pt(10)


 
    # 결과 저장
    prs.save(report_path)
    
    
if __name__ == "__main__":
    # 엑셀 파일 경로
    excel_file_path = '/home/kic/yskids/service/data/SAMSUNG'
    # 엑셀 파일을 활용하여 ppt 생성
    create_ppt_from_excel(excel_file_path)

